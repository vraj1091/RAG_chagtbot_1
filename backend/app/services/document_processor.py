"""Document processing service for extracting text from various file formats."""
import os
import io
import logging
from pathlib import Path
from typing import Optional, List, Tuple
import uuid

from app.core.config import settings

logger = logging.getLogger(__name__)

# Lazy imports for document processors - only imported when actually used
def _import_pdf_reader():
    from PyPDF2 import PdfReader
    return PdfReader

def _import_docx():
    from docx import Document as DocxDocument
    return DocxDocument

def _import_openpyxl():
    from openpyxl import load_workbook
    return load_workbook

def _import_pptx():
    from pptx import Presentation
    return Presentation

def _import_pil():
    from PIL import Image
    return Image

# OCR support
try:
    import pytesseract
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False

try:
    from pdf2image import convert_from_path, convert_from_bytes
    PDF2IMAGE_AVAILABLE = True
except ImportError:
    PDF2IMAGE_AVAILABLE = False


class DocumentProcessor:
    """Service for processing various document types."""
    
    SUPPORTED_EXTENSIONS = {
        'txt': 'text',
        'md': 'text',
        'csv': 'text',
        'pdf': 'pdf',
        'docx': 'docx',
        'doc': 'docx',
        'xlsx': 'xlsx',
        'xls': 'xlsx',
        'pptx': 'pptx',
        'png': 'image',
        'jpg': 'image',
        'jpeg': 'image',
        'bmp': 'image',
        'gif': 'image',
        'tiff': 'image',
        'tif': 'image'
    }
    
    @staticmethod
    def get_file_type(filename: str) -> str:
        """Get the processing type for a file."""
        ext = filename.lower().split('.')[-1]
        return DocumentProcessor.SUPPORTED_EXTENSIONS.get(ext, 'other')
    
    @staticmethod
    def is_supported(filename: str) -> bool:
        """Check if a file type is supported."""
        ext = filename.lower().split('.')[-1]
        return ext in DocumentProcessor.SUPPORTED_EXTENSIONS
    
    @staticmethod
    def generate_unique_filename(original_filename: str) -> str:
        """Generate a unique filename to avoid collisions."""
        ext = original_filename.split('.')[-1]
        unique_id = str(uuid.uuid4())[:8]
        safe_name = "".join(c for c in original_filename.rsplit('.', 1)[0] if c.isalnum() or c in '._- ')[:50]
        return f"{safe_name}_{unique_id}.{ext}"
    
    @staticmethod
    async def save_uploaded_file(file_content: bytes, filename: str, user_id: int) -> Tuple[str, str]:
        """Save an uploaded file and return (file_path, unique_filename)."""
        # Create user-specific upload directory
        user_upload_dir = settings.upload_path / str(user_id)
        user_upload_dir.mkdir(parents=True, exist_ok=True)
        
        # Generate unique filename
        unique_filename = DocumentProcessor.generate_unique_filename(filename)
        file_path = user_upload_dir / unique_filename
        
        # Save file
        with open(file_path, 'wb') as f:
            f.write(file_content)
        
        return str(file_path), unique_filename
    
    @staticmethod
    def extract_text(file_path: str, file_type: str) -> str:
        """Extract text from a document based on its type."""
        try:
            if file_type == 'text':
                return DocumentProcessor._extract_text_file(file_path)
            elif file_type == 'pdf':
                return DocumentProcessor._extract_pdf(file_path)
            elif file_type == 'docx':
                return DocumentProcessor._extract_docx(file_path)
            elif file_type == 'xlsx':
                return DocumentProcessor._extract_xlsx(file_path)
            elif file_type == 'pptx':
                return DocumentProcessor._extract_pptx(file_path)
            elif file_type == 'image':
                return DocumentProcessor._extract_image_text(file_path)
            else:
                logger.warning(f"Unsupported file type: {file_type}")
                return ""
        except Exception as e:
            logger.error(f"Error extracting text from {file_path}: {str(e)}")
            raise
    
    @staticmethod
    def _extract_text_file(file_path: str) -> str:
        """Extract text from a plain text file."""
        encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue
        
        # Fallback: read as binary and decode with errors='ignore'
        with open(file_path, 'rb') as f:
            return f.read().decode('utf-8', errors='ignore')
    
    @staticmethod
    def _extract_pdf(file_path: str) -> str:
        """Extract text from a PDF file with OCR fallback."""
        text_parts = []
        
        try:
            # Try regular text extraction first
            PdfReader = _import_pdf_reader()
            with open(file_path, 'rb') as f:
                reader = PdfReader(f)
                for page in reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(page_text)
            
            full_text = "\n".join(text_parts)
            
            # If no text extracted and OCR is available, try OCR
            if not full_text.strip() and OCR_AVAILABLE and PDF2IMAGE_AVAILABLE:
                logger.info(f"No text found in PDF, attempting OCR: {file_path}")
                full_text = DocumentProcessor._ocr_pdf(file_path)
            
            return full_text
            
        except Exception as e:
            logger.error(f"Error extracting PDF text: {str(e)}")
            # Try OCR as fallback
            if OCR_AVAILABLE and PDF2IMAGE_AVAILABLE:
                try:
                    return DocumentProcessor._ocr_pdf(file_path)
                except Exception as ocr_error:
                    logger.error(f"OCR fallback failed: {str(ocr_error)}")
            raise
    
    @staticmethod
    def _ocr_pdf(file_path: str) -> str:
        """Perform OCR on a PDF file."""
        text_parts = []
        
        try:
            images = convert_from_path(file_path)
            for i, image in enumerate(images):
                page_text = pytesseract.image_to_string(image)
                if page_text.strip():
                    text_parts.append(f"[Page {i+1}]\n{page_text}")
        except Exception as e:
            logger.error(f"OCR PDF error: {str(e)}")
            # Try with bytes if path doesn't work
            with open(file_path, 'rb') as f:
                images = convert_from_bytes(f.read())
                for i, image in enumerate(images):
                    page_text = pytesseract.image_to_string(image)
                    if page_text.strip():
                        text_parts.append(f"[Page {i+1}]\n{page_text}")
        
        return "\n\n".join(text_parts)
    
    @staticmethod
    def _extract_docx(file_path: str) -> str:
        """Extract text from a Word document."""
        DocxDocument = _import_docx()
        doc = DocxDocument(file_path)
        text_parts = []
        
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        
        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    text_parts.append(row_text)
        
        return "\n".join(text_parts)
    
    @staticmethod
    def _extract_xlsx(file_path: str) -> str:
        """Extract text from an Excel file."""
        load_workbook = _import_openpyxl()
        wb = load_workbook(file_path, data_only=True)
        text_parts = []
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            text_parts.append(f"[Sheet: {sheet_name}]")
            
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join(str(cell) for cell in row if cell is not None)
                if row_text.strip():
                    text_parts.append(row_text)
            
            text_parts.append("")  # Add blank line between sheets
        
        return "\n".join(text_parts)
    
    @staticmethod
    def _extract_pptx(file_path: str) -> str:
        """Extract text from a PowerPoint file."""
        Presentation = _import_pptx()
        prs = Presentation(file_path)
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"[Slide {slide_num}]"]
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            
            if len(slide_text) > 1:  # Only add if there's content besides the header
                text_parts.extend(slide_text)
                text_parts.append("")  # Blank line between slides
        
        return "\n".join(text_parts)
    
    @staticmethod
    def _extract_image_text(file_path: str) -> str:
        """Extract text from an image using OCR."""
        if not OCR_AVAILABLE:
            logger.warning("OCR not available. Install pytesseract for image text extraction.")
            return "[Image file - OCR not available]"
        
        try:
            Image = _import_pil()
            image = Image.open(file_path)
            text = pytesseract.image_to_string(image)
            return text.strip() if text.strip() else "[No text detected in image]"
        except Exception as e:
            logger.error(f"Error performing OCR on image: {str(e)}")
            raise
    
    @staticmethod
    def delete_file(file_path: str) -> bool:
        """Delete a file from the filesystem."""
        try:
            if os.path.exists(file_path):
                os.remove(file_path)
                return True
            return False
        except Exception as e:
            logger.error(f"Error deleting file {file_path}: {str(e)}")
            return False
